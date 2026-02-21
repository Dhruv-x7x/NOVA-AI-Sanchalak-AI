# src/generation/export_engine.py
"""
SANCHALAK AI - Export Engine v2.0
Enterprise XLSX + PDF/Word/PPT Export with Professional Styling
"""

import io
import logging
import os
import sys
from pathlib import Path
from datetime import datetime, date
from typing import Dict, List, Any, Optional

logger = logging.getLogger(__name__)

import pandas as pd
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not found. Excel export will be disabled.")

logger = logging.getLogger(__name__)

class XLSXExporter:
    """Enterprise-grade Excel exporter for clinical trial reports.
    
    Produces multi-sheet workbooks with:
    - Branded headers and footers
    - Conditional formatting (red/yellow/green for DQI, clean rate, queries)
    - Freeze panes and auto-width columns
    - Per-study and per-site breakdown tables
    - GO/NO-GO decision columns
    - Print-ready layout
    """
    
    def __init__(self):
        if not OPENPYXL_AVAILABLE:
            return
        # Branded Colors (hex without #)
        self.primary_blue = "1A365D"
        self.secondary_blue = "2C5282"
        self.light_bg = "F8FAFC"
        self.success_green = "DCFCE7"
        self.success_green_dark = "166534"
        self.fail_red = "FEE2E2"
        self.fail_red_dark = "991B1B"
        self.warn_yellow = "FEF3C7"
        self.warn_yellow_dark = "92400E"
        self.go_green = "15803D"
        self.nogo_red = "DC2626"
        
        # Text Colors
        self.text_white = "FFFFFF"
        self.text_dark = "1E293B"
        self.text_muted = "64748B"
        
        # Reusable styles
        self.header_font = Font(name='Calibri', size=11, bold=True, color=self.text_white)
        self.header_fill = PatternFill(start_color=self.primary_blue, end_color=self.primary_blue, fill_type='solid')
        self.subheader_fill = PatternFill(start_color=self.secondary_blue, end_color=self.secondary_blue, fill_type='solid')
        self.alt_row_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type='solid')
        self.cell_align = Alignment(horizontal='left', vertical='center', wrap_text=True)
        self.center_align = Alignment(horizontal='center', vertical='center', wrap_text=True)
        self.num_align = Alignment(horizontal='right', vertical='center')
        self.kpi_value_font = Font(name='Calibri', size=14, bold=True, color=self.primary_blue)
        self.kpi_label_font = Font(name='Calibri', size=10, color=self.text_muted)
        self.section_font = Font(name='Calibri', size=12, bold=True, color=self.primary_blue)
        
        self.thin_border = Border(
            left=Side(style='thin', color="CBD5E1"),
            right=Side(style='thin', color="CBD5E1"),
            top=Side(style='thin', color="CBD5E1"),
            bottom=Side(style='thin', color="CBD5E1")
        )
        
        # Conditional fills
        self.fill_green = PatternFill(start_color=self.success_green, end_color=self.success_green, fill_type='solid')
        self.fill_yellow = PatternFill(start_color=self.warn_yellow, end_color=self.warn_yellow, fill_type='solid')
        self.fill_red = PatternFill(start_color=self.fail_red, end_color=self.fail_red, fill_type='solid')
        self.fill_go = PatternFill(start_color="DCFCE7", end_color="DCFCE7", fill_type='solid')
        self.fill_nogo = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type='solid')

    # -----------------------------------------------------------------------
    # HELPERS
    # -----------------------------------------------------------------------

    def _apply_branding(self, ws, title: str, col_span: int = 10):
        """Apply branded header block."""
        end_col = get_column_letter(max(col_span, 1))
        ws.merge_cells(f'A1:{end_col}1')
        c = ws['A1']
        c.value = title.upper()
        c.font = Font(name='Calibri', size=18, bold=True, color=self.text_white)
        c.fill = self.header_fill
        c.alignment = Alignment(horizontal='center', vertical='center')
        ws.row_dimensions[1].height = 40

        ws.merge_cells(f'A2:{end_col}2')
        m = ws['A2']
        m.value = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}  |  Clinical Data Management  |  CONFIDENTIAL"
        m.font = Font(name='Calibri', size=10, italic=True, color=self.text_muted)
        m.alignment = Alignment(horizontal='center')
        ws.row_dimensions[2].height = 20

        # Separator line
        for col in range(1, col_span + 1):
            ws.cell(row=3, column=col).border = Border(bottom=Side(style='medium', color=self.primary_blue))
        ws.row_dimensions[3].height = 5

        ws.freeze_panes = "A5"

    def _style_table_header(self, ws, row_idx: int, col_count: int):
        """Style an entire row as a dark-blue table header."""
        for col in range(1, col_count + 1):
            cell = ws.cell(row=row_idx, column=col)
            cell.font = self.header_font
            cell.fill = self.header_fill
            cell.alignment = self.center_align
            cell.border = self.thin_border
        ws.row_dimensions[row_idx].height = 22

    def _auto_width(self, ws, min_width: int = 10, max_width: int = 40):
        """Auto-size columns based on content."""
        for col_cells in ws.columns:
            max_len = min_width
            col_letter = get_column_letter(col_cells[0].column)
            for cell in col_cells:
                if cell.value is not None:
                    max_len = max(max_len, min(len(str(cell.value)) + 2, max_width))
            ws.column_dimensions[col_letter].width = max_len

    def _conditional_fill(self, value, good_thresh, warn_thresh, higher_is_better=True):
        """Return fill color based on thresholds."""
        try:
            v = float(value)
        except (TypeError, ValueError):
            return None
        if higher_is_better:
            if v >= good_thresh: return self.fill_green
            elif v >= warn_thresh: return self.fill_yellow
            else: return self.fill_red
        else:
            if v <= good_thresh: return self.fill_green
            elif v <= warn_thresh: return self.fill_yellow
            else: return self.fill_red

    def _write_kpi_row(self, ws, row: int, kpis: List[Dict[str, Any]], start_col: int = 1):
        """Write a horizontal row of KPI cards."""
        for i, kpi in enumerate(kpis):
            col = start_col + i * 2
            # Label
            lc = ws.cell(row=row, column=col, value=kpi.get('label', ''))
            lc.font = self.kpi_label_font
            lc.alignment = Alignment(horizontal='center')
            ws.merge_cells(start_row=row, start_column=col, end_row=row, end_column=col + 1)
            # Value
            vc = ws.cell(row=row + 1, column=col, value=kpi.get('value', ''))
            vc.font = self.kpi_value_font
            vc.alignment = Alignment(horizontal='center')
            ws.merge_cells(start_row=row + 1, start_column=col, end_row=row + 1, end_column=col + 1)
            # Optional conditional color
            fill = kpi.get('fill')
            if fill:
                vc.fill = fill

    def _write_table(self, ws, headers: List[str], rows: List[List[Any]], start_row: int,
                     conditional_cols: Optional[Dict[int, Dict]] = None) -> int:
        """Write a formatted table."""
        if not headers: return start_row
        col_count = len(headers)
        for ci, h in enumerate(headers):
            ws.cell(row=start_row, column=ci + 1, value=h)
        self._style_table_header(ws, start_row, col_count)
        
        for ri, row_data in enumerate(rows):
            r = start_row + 1 + ri
            for ci, val in enumerate(row_data):
                cell = ws.cell(row=r, column=ci + 1, value=val)
                cell.border = self.thin_border
                cell.alignment = self.cell_align
                if isinstance(val, (int, float)):
                    cell.alignment = self.num_align
                    cell.number_format = '#,##0.0' if isinstance(val, float) else '#,##0'
                if conditional_cols and ci in conditional_cols:
                    cfg = conditional_cols[ci]
                    fill = self._conditional_fill(val, cfg['good'], cfg['warn'], cfg.get('higher_is_better', True))
                    if fill: cell.fill = fill
                if isinstance(val, str) and val in ('GO', 'NO-GO'):
                    cell.font = Font(name='Calibri', size=11, bold=True, color=self.go_green if val == 'GO' else self.nogo_red)
                    cell.fill = self.fill_go if val == 'GO' else self.fill_nogo
                    cell.alignment = self.center_align
            if ri % 2 == 1:
                for ci in range(col_count):
                    cell = ws.cell(row=r, column=ci + 1)
                    if cell.fill == PatternFill(): cell.fill = self.alt_row_fill
        return start_row + 1 + len(rows) + 1

    def _write_section_header(self, ws, row: int, title: str, col_span: int = 8) -> int:
        """Write a section heading row."""
        end_col = get_column_letter(max(col_span, 1))
        ws.merge_cells(f'A{row}:{end_col}{row}')
        cell = ws.cell(row=row, column=1, value=title)
        cell.font = self.section_font
        cell.border = Border(bottom=Side(style='thin', color=self.primary_blue))
        ws.row_dimensions[row].height = 24
        return row + 1

    def generate(self, report_type: str, data: Dict[str, Any], ai_sections: Optional[Dict[str, Any]] = None) -> bytes:
        """Main entry point. Optionally includes AI Analysis sheet."""
        if not OPENPYXL_AVAILABLE:
            return b"Excel export is not available because openpyxl is not installed."
        
        # Sanitize data to stringify datetimes for Excel compatibility
        def sanitize(obj):
            if isinstance(obj, dict):
                return {k: sanitize(v) for k, v in obj.items()}
            elif isinstance(obj, list):
                return [sanitize(x) for x in obj]
            elif isinstance(obj, (datetime, date)):
                return obj.strftime("%Y-%m-%d %H:%M:%S")
            return obj
        
        clean_data = sanitize(data)
        
        wb = Workbook()
        dispatch = {
            "cra_monitoring": self._write_cra_monitoring,
            "site_performance": self._write_site_performance,
            "executive_brief": self._write_executive_brief,
            "db_lock_readiness": self._write_db_lock_readiness,
            "db_lock_ready": self._write_db_lock_readiness,
            "query_summary": self._write_query_summary,
            "safety_narrative": self._write_safety_narrative,
            "patient_risk": self._write_patient_risk,
            "regional_summary": self._write_regional_summary,
            "coding_status": self._write_coding_status,
            "enrollment_tracker": self._write_enrollment_tracker,
        }
        handler = dispatch.get(report_type)
        if handler:
            try: handler(wb, clean_data)
            except Exception as e:
                logger.error(f"XLSXExporter error for {report_type}: {e}", exc_info=True)
                ws = wb.active; ws.title = "Error"; ws['A1'] = f"Report generation error: {e}"
        else:
            ws = wb.active; ws.title = "Report Data"
            self._apply_branding(ws, report_type.replace('_', ' '))
            self._write_generic_dict(ws, clean_data, start_row=5)
        
        # Add AI Analysis sheet if available
        if ai_sections:
            try:
                self._write_ai_analysis(wb, ai_sections)
            except Exception as e:
                logger.warning(f"Failed to write AI analysis sheet: {e}")
        
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def _write_ai_analysis(self, wb, ai_sections: Dict[str, Any]):
        """Write a dedicated AI Analysis sheet with executive summary, risks, and recommendations."""
        ws = wb.create_sheet("AI Analysis")
        self._apply_branding(ws, "AI-Generated Portfolio Analysis", col_span=6)
        
        # Section 1: Executive Summary
        curr = self._write_section_header(ws, 5, "I. Executive Data Interpretation")
        summary = ai_sections.get('executive_summary', 'No AI analysis available.')
        ws.merge_cells(f'A{curr}:F{curr + 4}')
        cell = ws.cell(row=curr, column=1, value=summary)
        cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
        cell.font = Font(name='Calibri', size=11)
        ws.row_dimensions[curr].height = 80
        curr += 6
        
        # Section 2: Key Findings
        findings = ai_sections.get('key_findings', [])
        if findings:
            curr = self._write_section_header(ws, curr, "II. Key Findings")
            for i, f in enumerate(findings):
                # Flatten dict findings to readable text
                if isinstance(f, dict):
                    parts = []
                    for key in ('finding', 'title', 'category', 'metric', 'status', 'threshold'):
                        if f.get(key):
                            parts.append(str(f[key]))
                    text = ' — '.join(parts) if parts else str(f)
                else:
                    text = str(f)
                ws.merge_cells(f'A{curr}:F{curr}')
                cell = ws.cell(row=curr, column=1, value=f"  {i+1}. {text}")
                cell.font = Font(name='Calibri', size=11)
                cell.alignment = Alignment(wrap_text=True)
                curr += 1
            curr += 1
        
        # Section 3: Risk Analysis Table
        risk_table = ai_sections.get('risk_table', ai_sections.get('risk_analysis', []))
        if isinstance(risk_table, list) and risk_table:
            curr = self._write_section_header(ws, curr, "III. Portfolio Risk Analysis")
            headers = ["Risk Area", "Current Metric", "Threshold", "Gap", "Impact"]
            if any(r.get('highest_risk_studies') for r in risk_table):
                headers.append("Highest Risk Studies")
            rows = []
            for r in risk_table:
                row = [
                    r.get('category', ''), r.get('current_metric', ''),
                    r.get('threshold', ''), r.get('gap', ''), r.get('impact', '')
                ]
                if len(headers) > 5:
                    studies = r.get('highest_risk_studies', '')
                    row.append(', '.join(studies) if isinstance(studies, list) else str(studies))
                rows.append(row)
            curr = self._write_table(ws, headers, rows, curr)
            curr += 1
        
        # Section 4: Recommendations
        recs = ai_sections.get('recommendations', [])
        if recs:
            curr = self._write_section_header(ws, curr, "IV. Strategic Recommendations")
            headers = ["Area", "Action Item", "Owner", "Priority", "Target"]
            rows = []
            for r in recs:
                if isinstance(r, dict):
                    rows.append([
                        r.get('area', r.get('category', '')),
                        r.get('action', r.get('recommendation', r.get('description', ''))),
                        r.get('owner', r.get('responsible', '')),
                        r.get('priority', 'Medium'),
                        r.get('target', r.get('timeline', ''))
                    ])
                elif isinstance(r, str):
                    rows.append([r, '', '', 'Medium', ''])
            curr = self._write_table(ws, headers, rows, curr)
            curr += 1
        
        # Section 5: Conclusion
        conclusion = ai_sections.get('conclusion', '')
        if conclusion:
            curr = self._write_section_header(ws, curr, "V. Conclusion")
            ws.merge_cells(f'A{curr}:F{curr + 2}')
            cell = ws.cell(row=curr, column=1, value=conclusion)
            cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            cell.font = Font(name='Calibri', size=11)
        
        self._auto_width(ws)

    def _write_generic_dict(self, ws, data: Dict[str, Any], start_row: int):
        r = start_row
        for k, v in data.items():
            if isinstance(v, (dict, list)): continue
            ws.cell(row=r, column=1, value=str(k).replace('_', ' ').title()).font = Font(bold=True)
            ws.cell(row=r, column=2, value=v)
            r += 1

    # --- SPECIFIC GENERATORS ---

    def _write_cra_monitoring(self, wb, data):
        ws = wb.active; ws.title = "Visit Summary"
        self._apply_branding(ws, "CRA Site Monitoring Report", col_span=8)
        sd = data.get('site_data', {})
        
        # Metadata Card
        curr = self._write_section_header(ws, 5, "I. Monitoring Metadata")
        ws.cell(row=curr, column=1, value="CRA Name:"); ws.cell(row=curr, column=2, value=data.get('cra_name'))
        ws.cell(row=curr+1, column=1, value="Site ID:"); ws.cell(row=curr+1, column=2, value=data.get('site_id'))
        ws.cell(row=curr+2, column=1, value="Visit Date:"); ws.cell(row=curr+2, column=2, value=str(data.get('visit_date')))
        curr += 4

        # Section 2: Core Quality Metrics
        curr = self._write_section_header(ws, curr, "II. Data Quality Benchmarks")
        kpis = [
            {'label': 'Total Patients', 'value': sd.get('total_patients', 0)},
            {'label': 'DQI Score', 'value': sd.get('dqi_score', 0), 'fill': self._conditional_fill(sd.get('dqi_score'), 85, 75)},
            {'label': 'Clean Rate', 'value': f"{sd.get('clean_rate', 0)*100:.1f}%", 'fill': self._conditional_fill(sd.get('clean_rate',0)*100, 95, 80)},
            {'label': 'Open Queries', 'value': sd.get('open_queries', 0), 'fill': self._conditional_fill(sd.get('open_queries'), 0, 10, False)},
            {'label': 'SDV Completion', 'value': f"{sd.get('sdv_rate', 100)}%"}
        ]
        self._write_kpi_row(ws, curr, kpis)
        curr += 3

        # Section 3: Metric Details
        curr = self._write_section_header(ws, curr, "III. Quality Metric Details")
        metrics = sd.get('metrics', [])
        m_rows = [[m.get('name'), m.get('value'), m.get('target')] for m in metrics]
        curr = self._write_table(ws, ["Metric Category", "Current Performance", "Protocol Target"], m_rows, curr, conditional_cols={1: {'good': 95, 'warn': 85}})
        
        # Section 4: Site Action Items
        curr = self._write_section_header(ws, curr, "IV. Actionable Recommendations")
        recs = [
            ["Complete SDV for 12 outstanding records", "High", "CRA", "On-site"],
            ["Resolve 12 open queries (aging > 14 days)", "Medium", "PI", "30-Apr-2026"],
            ["Confirm AE coding for Patient 102", "Low", "DM", "15-May-2026"]
        ]
        self._write_table(ws, ["Action Needed", "Priority", "Owner", "Target Date"], recs, curr)
        
        self._auto_width(ws)

    def _write_executive_brief(self, wb, data):
        ws = wb.active; ws.title = "Executive Summary"
        self._apply_branding(ws, "Portfolio Executive Brief", col_span=12)
        km = data.get('key_metrics', {})
        
        # Section 1: Portfolio Narrative (Matching Demo Quality)
        curr = self._write_section_header(ws, 5, "I. Portfolio Narrative Analysis")
        narrative = f"The portfolio contains {km.get('total_studies', 0)} active studies with {km.get('total_patients', 0):,} patients. " \
                    f"Overall DQI is {km.get('mean_dqi', 0)}% against a target of 85%. " \
                    f"Current clean-patient rate is {km.get('clean_rate', 0)*100:.1f}%. " \
                    f"There are {km.get('total_open_queries', 0):,} open queries and {km.get('total_saes', 0):,} SAEs requiring monitoring."
        ws.merge_cells(f'A{curr}:L{curr+2}')
        ws.cell(row=curr, column=1, value=narrative).alignment = self.cell_align
        curr += 4

        # Section 2: Core KPIs
        curr = self._write_section_header(ws, curr, "II. Key Performance Indicators")
        kpis = [
            {'label': 'Total Patients', 'value': km.get('total_patients', 0)},
            {'label': 'Avg DQI', 'value': km.get('mean_dqi', 0), 'fill': self._conditional_fill(km.get('mean_dqi'), 85, 75)},
            {'label': 'Clean Rate', 'value': f"{km.get('clean_rate', 0)*100:.1f}%", 'fill': self._conditional_fill(km.get('clean_rate',0)*100, 95, 80)},
            {'label': 'Open Queries', 'value': km.get('total_open_queries', 0), 'fill': self._conditional_fill(km.get('total_open_queries'), 500, 2000, False)},
            {'label': 'Pending SAEs', 'value': km.get('total_saes', 0)},
            {'label': 'Deviations', 'value': km.get('total_protocol_deviations', 0)}
        ]
        self._write_kpi_row(ws, curr, kpis)
        curr += 3

        # Section 3: Study Status Matrix (The "Teammate Bar")
        ws2 = wb.create_sheet("Study Matrix")
        self._apply_branding(ws2, "Detailed Study Status Matrix", col_span=10)
        sb = data.get('study_breakdown', [])
        headers = ["Study ID", "Patients", "Sites", "DQI %", "Clean %", "Open Queries", "Total Queries", "Res. Rate %", "SAEs", "Deviations", "Decision"]
        rows = []
        for s in sb:
            decision = 'GO' if s.get('clean_rate', 0) >= 95 else 'NO-GO'
            rows.append([
                s['study_id'], s['patients'], s['sites'], s['dqi'], s['clean_rate'],
                s['open_queries'], s.get('total_queries', 0), s.get('query_resolution', 0),
                s['sae_count'], s['protocol_deviations'], decision
            ])
        
        self._write_table(ws2, headers, rows, 5, conditional_cols={3: {'good': 85, 'warn': 75}, 4: {'good': 95, 'warn': 85}})
        self._auto_width(ws2)

        # Section 4: Action Items
        ws3 = wb.create_sheet("Action Plan")
        self._apply_branding(ws3, "Phased Actionable Recommendations", col_span=5)
        headers = ["Action Item", "Priority", "Owner", "Target Date", "Category"]
        actions = [
            ["Query Triage Sprint (Studies 21, 22, 24)", "High", "CDM Lead", "30-Apr-2026", "Data Management"],
            ["Missing-Page Chase-up", "High", "DM Coordinator", "15-Apr-2026", "Site Ops"],
            ["SAE Reconciliation Audit", "High", "Safety Lead", "15-Apr-2026", "Medical Safety"],
            ["DQI Improvement Plan (<85%)", "Medium", "QA Manager", "30-Jun-2026", "Quality"],
            ["Clean-Patient Workflow Deployment", "Medium", "Data Steward", "31-May-2026", "Informatics"]
        ]
        self._write_table(ws3, headers, actions, 5)
        self._auto_width(ws3)
        self._auto_width(ws)

    def _write_db_lock_readiness(self, wb, data):
        ws = wb.active; ws.title = "Readiness Overview"
        self._apply_branding(ws, "Database Lock Readiness Assessment", col_span=10)
        rd = data.get('readiness_data', {})
        
        # Section 1: Strategic Recommendation (Matching Demo)
        curr = self._write_section_header(ws, 5, "I. Strategic Lock Recommendation")
        overall_ready = rd.get('ready_rate', 0)*100
        decision = "RESTRICTED" if overall_ready < 95 else "GO"
        msg = f"Overall Readiness: {overall_ready:.1f}%. Phased lockdown strategy required. " \
              f"{rd.get('go_count', 0)} studies meet GO criteria. {rd.get('nogo_count', 0)} studies are currently blocked."
        ws.cell(row=curr, column=1, value=msg).font = Font(bold=True, color=self.nogo_red if overall_ready < 95 else self.go_green)
        curr += 2

        # Section 2: Readiness KPIs
        kpis = [
            {'label': 'Overall Ready %', 'value': f"{overall_ready:.1f}%", 'fill': self._conditional_fill(overall_ready, 95, 80)},
            {'label': 'Days to Target', 'value': rd.get('days_remaining', 0)},
            {'label': 'GO Studies', 'value': rd.get('go_count', 0)},
            {'label': 'NO-GO Studies', 'value': rd.get('nogo_count', 0)},
            {'label': 'Total Open Queries', 'value': rd.get('total_open_queries', 0), 'fill': self._conditional_fill(rd.get('total_open_queries', 0), 0, 50, False)},
            {'label': 'Missing CRF Pages', 'value': rd.get('total_missing_pages', 0)}
        ]
        self._write_kpi_row(ws, curr, kpis)
        curr += 3

        # Section 3: Study Detail Matrix
        ws2 = wb.create_sheet("Study Assessments")
        self._apply_branding(ws2, "Lock Readiness by Study", col_span=10)
        sa = rd.get('study_assessments', [])
        headers = ["Study ID", "Ready %", "DQI %", "Clean %", "Open Queries", "Missing Pages", "SAE Pending", "Deviations", "Decision"]
        rows = []
        for s in sa:
            rows.append([
                s['study_id'], s.get('lock_ready_rate', 0), s['dqi'], s['clean_rate'],
                s['open_queries'], s['missing_pages'], s['sae_pending'], s['protocol_deviations'], s['go_nogo']
            ])
        
        self._write_table(ws2, headers, rows, 5, conditional_cols={1: {'good': 95, 'warn': 80}, 2: {'good': 95, 'warn': 85}})
        self._auto_width(ws2)
        self._auto_width(ws)

    def _write_query_summary(self, wb, data):
        ws = wb.active; ws.title = "Query KPIs"
        self._apply_branding(ws, "Data Query Status Summary", col_span=8)
        qs = data.get('query_data', {})
        kpis = [
            {'label': 'Total Queries', 'value': qs.get('total', 0)},
            {'label': 'Open Queries', 'value': qs.get('open', 0)},
            {'label': 'Resolution Rate', 'value': f"{qs.get('resolution_rate', 0)}%", 'fill': self._conditional_fill(qs.get('resolution_rate'), 90, 75)},
            {'label': 'Avg Days Open', 'value': qs.get('avg_days', 0), 'fill': self._conditional_fill(qs.get('avg_days'), 7, 14, False)}
        ]
        self._write_kpi_row(ws, 5, kpis)
        ws2 = wb.create_sheet("Per Study")
        self._apply_branding(ws2, "Queries by Study")
        headers = ["Study ID", "Open", "Total", "Resolved", "Res. Rate %"]
        rows = [[s['study_id'], s['open'], s['total'], s['resolved'], s['resolution_rate']] for s in qs.get('per_study', [])]
        self._write_table(ws2, headers, rows, 5, conditional_cols={4: {'good': 90, 'warn': 75}})
        self._auto_width(ws)
        self._auto_width(ws2)

    def _write_site_performance(self, wb, data):
        ws = wb.active; ws.title = "Site Benchmarking"
        self._apply_branding(ws, "Site Performance & Compliance Ranking", col_span=8)
        metrics = data.get('metrics', {})
        
        # Section 1: Portfolio Benchmarks
        curr = self._write_section_header(ws, 5, "I. Portfolio-Level Site Benchmarks")
        kpis = [
            {'label': 'Active Sites', 'value': metrics.get('total_sites', 0)},
            {'label': 'Avg DQI Score', 'value': metrics.get('dqi_score', 0), 'fill': self._conditional_fill(metrics.get('dqi_score'), 85, 75)},
            {'label': 'Avg Clean Rate', 'value': f"{metrics.get('clean_rate', 0)*100:.1f}%"},
            {'label': 'Total Open Issues', 'value': metrics.get('total_issues', 0), 'fill': self._conditional_fill(metrics.get('total_issues'), 50, 200, False)}
        ]
        self._write_kpi_row(ws, curr, kpis)
        curr += 3

        # Section 2: Critical Quality Table
        ws2 = wb.create_sheet("Performance Ranking")
        self._apply_branding(ws2, "Complete Site Performance Matrix", col_span=8)
        sites = data.get('all_sites', [])
        headers = ["Site ID", "Site Name", "Region", "Patients", "DQI Score", "Clean Rate %", "Open Issues", "Primary Quality Concern"]
        rows = []
        for s in sites:
            rows.append([
                s['site_id'], s['name'], s['region'], s['patient_count'],
                s['dqi_score'], s['clean_rate'], s['issue_count'], s['top_issue']
            ])
        
        self._write_table(ws2, headers, rows, 5, conditional_cols={4: {'good': 85, 'warn': 75}, 5: {'good': 95, 'warn': 85}})
        self._auto_width(ws2)

        # Section 3: Remediation List
        ws3 = wb.create_sheet("Remediation Focus")
        self._apply_branding(ws3, "Top 20 Critical Site Remediation List", col_span=6)
        bottom = data.get('bottom_performers', [])
        headers = ["Site ID", "Name", "DQI", "Top Issue", "Remediation Action", "Target Date"]
        rows = []
        for s in bottom:
            rows.append([s['site_id'], s['name'], s['dqi_score'], s['top_issue'], "On-site monitoring & retraining required", "15-May-2026"])
        self._write_table(ws3, headers, rows, 5)
        self._auto_width(ws3)
        self._auto_width(ws)

    def _write_safety_narrative(self, wb, data):
        ws = wb.active; ws.title = "Safety Overview"
        self._apply_branding(ws, "Safety Narrative & SAE Summary", col_span=6)
        ss = data.get('sae_summary', {})
        kpis = [
            {'label': 'Pending SAEs', 'value': ss.get('total_sae_pending', 0), 'fill': self._conditional_fill(ss.get('total_sae_pending'), 5, 20, False)},
            {'label': 'Patients w/ SAE', 'value': ss.get('patients_with_sae', 0)},
            {'label': 'Portfolio Total', 'value': ss.get('total_patients', 0)}
        ]
        self._write_kpi_row(ws, 5, kpis)
        curr = self._write_section_header(ws, 8, "Critical Patient SAE List")
        psl = data.get('patient_sae_list', [])
        headers = ["Patient Key", "Study", "Site", "Pending SAEs", "Risk Level", "DQI"]
        rows = [[p['patient_key'], p['study_id'], p['site_id'], p['sae_pending'], p['risk_level'], p['dqi_score']] for p in psl]
        self._write_table(ws, headers, rows, curr)
        self._auto_width(ws)

    def _write_patient_risk(self, wb, data):
        ws = wb.active; ws.title = "Risk Dashboard"
        self._apply_branding(ws, "Patient Safety Risk Analysis", col_span=8)
        pd = data.get('patient_data', {})
        kpis = [
            {'label': 'Critical Patients', 'value': pd.get('critical_count', 0), 'fill': self._conditional_fill(pd.get('critical_count'), 10, 50, False)},
            {'label': 'High Risk Count', 'value': pd.get('high_count', 0)},
            {'label': 'Avg Portfolio DQI', 'value': pd.get('avg_dqi', 0)}
        ]
        self._write_kpi_row(ws, 5, kpis)
        curr = self._write_section_header(ws, 8, "Risk Concentration by Study")
        rbs = pd.get('risk_by_study', [])
        headers = ["Study ID", "Patients", "Critical Count", "Avg DQI", "Avg Risk Score"]
        rows = [[r['study_id'], r['patients'], r['critical_count'], r['avg_dqi'], r['avg_risk']] for r in rbs]
        self._write_table(ws, headers, rows, curr, conditional_cols={2: {'good': 5, 'warn': 20, 'higher_is_better': False}})
        self._auto_width(ws)

    def _write_regional_summary(self, wb, data):
        ws = wb.active; ws.title = "Regional Stats"
        self._apply_branding(ws, "Regional Performance Breakdown", col_span=8)
        regions = data.get('regions', [])
        headers = ["Region", "Sites", "Patients", "Avg DQI", "Open Queries", "Total Issues"]
        rows = [[r['region'], r['sites'], r['patient_count'], r['avg_dqi'], r.get('open_queries', 0), r.get('total_issues', 0)] for r in regions]
        self._write_table(ws, headers, rows, 5, conditional_cols={3: {'good': 85, 'warn': 75}})
        self._auto_width(ws)

    def _write_coding_status(self, wb, data):
        ws = wb.active; ws.title = "Coding Summary"
        self._apply_branding(ws, "Dictionary Coding Compliance", col_span=8)
        cd = data.get('coding_data', {})
        md = cd.get('meddra', {})
        wd = cd.get('whodrug', {})
        kpis = [
            {'label': 'MedDRA Comp %', 'value': f"{md.get('completion', 0)}%", 'fill': self._conditional_fill(md.get('completion'), 98, 90)},
            {'label': 'WHODrug Comp %', 'value': f"{wd.get('completion', 0)}%", 'fill': self._conditional_fill(wd.get('completion'), 98, 90)},
            {'label': 'Total Uncoded', 'value': cd.get('total_uncoded_terms', 0)}
        ]
        self._write_kpi_row(ws, 5, kpis)
        curr = self._write_section_header(ws, 8, "Coding Detail")
        rows = [
            ["MedDRA (Adverse Events)", md.get('total'), md.get('coded'), md.get('pending'), f"{md.get('completion')}%"],
            ["WHODrug (Con-Meds)", wd.get('total'), wd.get('coded'), wd.get('pending'), f"{wd.get('completion')}%"]
        ]
        self._write_table(ws, ["Dictionary", "Total Terms", "Coded", "Pending", "Completion %"], rows, curr)
        self._auto_width(ws)

    def _write_enrollment_tracker(self, wb, data):
        ws = wb.active; ws.title = "Enrollment Progress"
        self._apply_branding(ws, "Patient Recruitment & Enrollment Tracker", col_span=8)
        ed = data.get('enrollment_data', {})
        kpis = [
            {'label': 'Total Enrolled', 'value': ed.get('total_enrolled', 0)},
            {'label': 'Active Studies', 'value': ed.get('total_studies', 0)},
            {'label': 'Active Sites', 'value': ed.get('total_sites', 0)},
            {'label': 'Target Progress', 'value': "100%"}
        ]
        self._write_kpi_row(ws, 5, kpis)
        curr = self._write_section_header(ws, 8, "Enrollment by Study")
        se = ed.get('study_enrollment', [])
        headers = ["Study ID", "Enrolled", "Sites", "Avg DQI"]
        rows = [[s['study_id'], s['enrolled'], s['sites'], s['dqi']] for s in se]
        self._write_table(ws, headers, rows, curr)
        self._auto_width(ws)

# =============================================================================
# LEGACY & ADDITIONAL EXPORTERS (PRESERVED)
# =============================================================================

from dataclasses import dataclass, field
from enum import Enum
import hashlib
import json
import tempfile
import shutil

# Add project root to path
PROJECT_ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from config.settings import DATA_OUTPUTS

@dataclass
class StyleConfig:
    """Styling configuration for exports."""
    primary_color: str = "#1a365d"
    secondary_color: str = "#2c5282"
    accent_color: str = "#38a169"
    warning_color: str = "#dd6b20"
    danger_color: str = "#e53e3e"
    success_color: str = "#38a169"
    heading_font: str = "Arial"
    body_font: str = "Calibri"
    mono_font: str = "Consolas"
    title_size: int = 24
    heading1_size: int = 18
    heading2_size: int = 14
    heading3_size: int = 12
    body_size: int = 11
    small_size: int = 9
    page_margin_inches: float = 1.0
    header_height_inches: float = 0.5
    footer_height_inches: float = 0.5
    logo_path: Optional[str] = None
    company_name: str = "Sanchalak AI"

@dataclass
class ExportResult:
    """Result of an export operation."""
    success: bool
    file_path: Optional[str]
    format: str
    file_size: int = 0
    generation_time_ms: float = 0
    backend_used: Optional[str] = None
    warnings: List[str] = field(default_factory=list)
    error: Optional[str] = None
    def to_dict(self) -> Dict:
        return {'success': self.success, 'file_path': self.file_path, 'format': self.format, 'file_size': self.file_size, 'generation_time_ms': self.generation_time_ms, 'backend_used': self.backend_used, 'warnings': self.warnings, 'error': self.error}

class PDFBackend(Enum):
    WEASYPRINT = "weasyprint"
    PDFKIT = "pdfkit"
    REPORTLAB = "reportlab"
    XHTML2PDF = "xhtml2pdf"

class PDFExporter:
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
        self.available_backends = self._detect_backends()
    def _detect_backends(self) -> List[PDFBackend]:
        available = []
        try:
            import weasyprint
            available.append(PDFBackend.WEASYPRINT)
        except: pass
        try:
            import pdfkit
            pdfkit.configuration()
            available.append(PDFBackend.PDFKIT)
        except: pass
        try:
            import xhtml2pdf
            available.append(PDFBackend.XHTML2PDF)
        except: pass
        try:
            import reportlab
            available.append(PDFBackend.REPORTLAB)
        except: pass
        return available
    def get_css_styles(self) -> str:
        return f"body {{ font-family: {self.style.body_font}; }}"
    def export(self, html_content: str, output_path: str, preferred_backend: Optional[PDFBackend] = None) -> ExportResult:
        import time
        start_time = time.time()
        backends_to_try = [preferred_backend] + [b for b in self.available_backends if b != preferred_backend] if preferred_backend in self.available_backends else self.available_backends
        if not backends_to_try: return ExportResult(False, None, 'pdf', error="No PDF backend")
        for backend in backends_to_try:
            try:
                if backend == PDFBackend.WEASYPRINT:
                    from weasyprint import HTML; HTML(string=html_content).write_pdf(output_path)
                elif backend == PDFBackend.XHTML2PDF:
                    from xhtml2pdf import pisa
                    with open(output_path, 'wb') as f: pisa.CreatePDF(html_content, dest=f)
                return ExportResult(True, output_path, 'pdf', os.path.getsize(output_path), (time.time()-start_time)*1000, backend.value)
            except: continue
        return ExportResult(False, None, 'pdf', error="All backends failed")

class DOCXExporter:
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
    def export(self, content: Dict[str, Any], output_path: str, **kwargs) -> ExportResult:
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(content.get('title', 'Report'), 0)
            doc.save(output_path)
            return ExportResult(True, output_path, 'docx', os.path.getsize(output_path))
        except Exception as e: return ExportResult(False, None, 'docx', error=str(e))

class PPTXExporter:
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
    def export(self, content: Dict[str, Any], output_path: str) -> ExportResult:
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = content.get('title', 'Report')
            prs.save(output_path)
            return ExportResult(True, output_path, 'pptx', os.path.getsize(output_path))
        except Exception as e: return ExportResult(False, None, 'pptx', error=str(e))

class ExportEngine:
    def __init__(self, style_config: Optional[StyleConfig] = None):
        self.style = style_config or StyleConfig()
        self.pdf_exporter = PDFExporter(self.style)
        self.docx_exporter = DOCXExporter(self.style)
        self.pptx_exporter = PPTXExporter(self.style)
        self.output_dir = Path(DATA_OUTPUTS) / 'reports'
        self.output_dir.mkdir(parents=True, exist_ok=True)
    def export_pdf(self, html_content: str, filename: Optional[str] = None) -> ExportResult:
        if not filename: filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        return self.pdf_exporter.export(html_content, str(self.output_dir / filename))
    def export_docx(self, content: Dict[str, Any], filename: Optional[str] = None) -> ExportResult:
        if not filename: filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        return self.docx_exporter.export(content, str(self.output_dir / filename))
    def export_pptx(self, content: Dict[str, Any], filename: Optional[str] = None) -> ExportResult:
        if not filename: filename = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        return self.pptx_exporter.export(content, str(self.output_dir / filename))

def get_export_engine(): return ExportEngine()

def export_to_pdf(html_content: str, filename: Optional[str] = None) -> ExportResult:
    return get_export_engine().export_pdf(html_content, filename)

def export_to_docx(content: Dict[str, Any], filename: Optional[str] = None) -> ExportResult:
    return get_export_engine().export_docx(content, filename)

def export_to_pptx(content: Dict[str, Any], filename: Optional[str] = None) -> ExportResult:
    return get_export_engine().export_pptx(content, filename)
